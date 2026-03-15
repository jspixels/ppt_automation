# Imports
import json

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

import google.generativeai as genai
from google import genai
from google.genai import types


# Extract text shapes from PPT
def extract_shapes(ppt_path):

    prs = Presentation(ppt_path)

    all_slides = []

    for slide in prs.slides:

        shapes = []

        for shape in slide.shapes:

            if shape.has_text_frame:

                text = shape.text.strip()

                if text:
                    shapes.append({
                        "text": text,
                        "x": shape.left,
                        "y": shape.top
                    })

        all_slides.append(shapes)

    return all_slides


# Reconstruct slide grid using shape coordinates
def reconstruct_grid(shapes):

    """
    shapes = list of dicts:
    {
        "text": "...",
        "x": shape.left,
        "y": shape.top
    }
    """

    if not shapes:
        return []

    # Row grouping tolerance
    ROW_THRESHOLD = 200000

    rows = []

    for shape in shapes:

        placed = False

        for row in rows:

            if abs(row["y"] - shape["y"]) < ROW_THRESHOLD:
                row["cells"].append(shape)
                placed = True
                break

        if not placed:
            rows.append({
                "y": shape["y"],
                "cells": [shape]
            })

    # Sort rows vertically
    rows = sorted(rows, key=lambda r: r["y"])

    grid = []

    for row in rows:

        # Sort cells horizontally
        cells = sorted(row["cells"], key=lambda c: c["x"])

        grid.append([c["text"] for c in cells])

    return grid


# Convert grid into slide JSON structure
def build_slide_json(grid):

    if not grid:
        return None

    slide_json = {
        "title": None,
        "tables": [],
        "text_blocks": [],
        "layout_features": {},
        "slide_type": None
    }

    # Detect title
    if len(grid[0]) == 1:
        slide_json["title"] = grid[0][0]
        grid = grid[1:]

    # Detect table slide
    if grid and len(grid[0]) >= 3:

        headers = grid[0]
        data_rows = grid[1:]

        # Detect real column count
        max_cols = max(len(r) for r in data_rows)
        header_cols = len(headers)

        # Fix missing header columns
        if max_cols > header_cols:
            missing = max_cols - header_cols
            headers = [""] * missing + headers

        # Normalize rows
        normalized_rows = []

        for row in data_rows:

            if len(row) < max_cols:
                row += [""] * (max_cols - len(row))

            normalized_rows.append(row)

        table = {
            "rows": len(normalized_rows) + 1,
            "columns": max_cols,
            "headers": headers,
            "data": normalized_rows
        }

        slide_json["tables"].append(table)
        slide_json["slide_type"] = "table"

    # Detect bullet slide
    else:

        slide_json["text_blocks"] = [cell for row in grid for cell in row]
        slide_json["slide_type"] = "bullet"

    # Layout metadata
    slide_json["layout_features"] = {
        "rows_detected": len(grid),
        "columns_detected": max(len(r) for r in grid)
    }

    return slide_json


# Refine slide JSON using Gemini
def refine_json_with_gemini(raw_input_json, api_key):

    client = genai.Client(api_key=api_key)

    system_prompt = """
    You are an expert PowerPoint presentation designer.

    Your task is to transform noisy raw slide JSON into a clean, professional structure
    suitable for generating a visually appealing presentation.

    Instructions:
    1. The input JSON may contain multiple slides.
    2. Clean all noise characters (\\x0b, \\xa0, extra spaces, broken sentences).
    3. Fix grammar and rewrite professional, and presentation-ready.
    4. Do NOT strictly follow the provided slide_type. Analyze the content and choose
       the best visual layout: "table", "bullet", or "paragraph".
    5. Determine the overall presentation theme from categories such as:
       business, corporate strategy, technical, education, research, startup pitch.
    6. Choose ONE consistent color palette for the entire presentation.
       Define theme_colors once at the top level — every slide must inherit it.
       theme_colors must include:
         - primary      : dominant background/header color  [R, G, B]
         - primary_text : text on primary background        [R, G, B]
         - accent       : highlight / accent color          [R, G, B]
         - slide_bg     : slide background color            [R, G, B]
         - body_text    : main body text color              [R, G, B]
    7. For bullet slides, each bullet must be a single, concise line (≤ 12 words).
    8. Make sure the color selection for theme appears soothing. 
    9. Slide background colour should be either white, dark grey or black. The table cell background colour should be different then the defined slide background colour.
    10. Include some spaces between each row in table.
    11. Include a design_score block explaining layout and theme choices.
    12. Output JSON must strictly follow the schema.
    """

    response_schema = {
        "type": "OBJECT",
        "properties": {
            "presentation_theme": {"type": "STRING"},
            "theme_colors": {
                "type": "OBJECT",
                "required": ["primary", "primary_text", "accent", "slide_bg", "body_text"],
                "properties": {
                    "primary": {"type": "ARRAY", "items": {"type": "INTEGER"}},
                    "primary_text": {"type": "ARRAY", "items": {"type": "INTEGER"}},
                    "accent": {"type": "ARRAY", "items": {"type": "INTEGER"}},
                    "slide_bg": {"type": "ARRAY", "items": {"type": "INTEGER"}},
                    "body_text": {"type": "ARRAY", "items": {"type": "INTEGER"}}
                }
            },
            "slides": {
                "type": "ARRAY",
                "items": {
                    "type": "OBJECT",
                    "properties": {
                        "slide_metadata": {
                            "type": "OBJECT",
                            "required": ["title", "layout_type"],
                            "properties": {
                                "title": {"type": "STRING"},
                                "layout_type": {
                                    "type": "STRING",
                                    "enum": ["table", "bullet", "paragraph"]
                                }
                            }
                        },
                        "table_data": {
                            "type": "OBJECT",
                            "nullable": True,
                            "properties": {
                                "rows": {"type": "INTEGER"},
                                "columns": {"type": "INTEGER"},
                                "headers": {"type": "ARRAY", "items": {"type": "STRING"}},
                                "content": {
                                    "type": "ARRAY",
                                    "items": {"type": "ARRAY", "items": {"type": "STRING"}}
                                }
                            }
                        },
                        "text_content": {
                            "type": "ARRAY",
                            "items": {"type": "STRING"}
                        },
                        "design_score": {
                            "type": "OBJECT",
                            "properties": {
                                "layout_reason": {"type": "STRING"},
                                "theme_reason": {"type": "STRING"},
                                "readability": {"type": "NUMBER"},
                                "visual_clarity": {"type": "NUMBER"}
                            }
                        }
                    }
                }
            }
        }
    }

    prompt = f"""
    Convert the following raw slide JSON into a refined professional presentation format.

    Input JSON:
    {json.dumps(raw_input_json)}

    Output should strictly follow the defined schema and contain clean, visually optimized slide data.
    """

    response = client.models.generate_content(
        model="gemini-2.0-flash",
        contents=prompt,
        config=types.GenerateContentConfig(
            system_instruction=system_prompt,
            response_mime_type="application/json",
            response_schema=response_schema,
            temperature=0.2
        )
    )

    return json.loads(response.text)


# Slide dimension constants
SLIDE_W = Inches(10)
SLIDE_H = Inches(7.5)

# Layout constants
MARGIN_X = Inches(0.5)
TITLE_TOP = Inches(0.25)
TITLE_H = Inches(0.75)
CONTENT_TOP = Inches(1.2)
CONTENT_H = Inches(5.8)
CONTENT_W = Inches(9.0)


# RGB helper
def _rgb(lst, fallback=(30, 30, 30)):
    try:
        return RGBColor(*lst)
    except Exception:
        return RGBColor(*fallback)


# Set slide background
def _set_slide_bg(slide, color_rgb):

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color_rgb


# Add title banner
def _add_title(slide, text, tc):

    primary = _rgb(tc["primary"])
    prim_text = _rgb(tc["primary_text"])

    title_box = slide.shapes.add_shape(
        1,
        MARGIN_X,
        TITLE_TOP,
        CONTENT_W,
        TITLE_H
    )

    title_box.fill.solid()
    title_box.fill.fore_color.rgb = primary
    title_box.line.fill.background()

    tf = title_box.text_frame
    tf.word_wrap = False
    tf.margin_left = Inches(0.15)
    tf.margin_right = Inches(0.15)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    p = tf.paragraphs[0]
    p.text = text
    p.alignment = PP_ALIGN.LEFT
    p.font.bold = True
    p.font.size = Pt(22)
    p.font.color.rgb = prim_text
    p.font.name = "Calibri"


# Add bullet slide
def _add_bullet_slide(slide, text_items, tc):

    accent = _rgb(tc["accent"])
    body_text = _rgb(tc["body_text"])

    txBox = slide.shapes.add_textbox(
        MARGIN_X,
        CONTENT_TOP,
        CONTENT_W,
        CONTENT_H
    )

    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    n = len(text_items)
    fsize = max(11, min(16, int(48 / max(n, 1))))

    for i, item in enumerate(text_items):

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        run = p.add_run()
        run.text = "● "
        run.font.color.rgb = accent
        run.font.size = Pt(fsize)
        run.font.name = "Calibri"

        run2 = p.add_run()
        run2.text = item.strip()
        run2.font.color.rgb = body_text
        run2.font.size = Pt(fsize)
        run2.font.name = "Calibri"

        p.space_before = Pt(4)
        p.space_after = Pt(4)
        p.level = 0


# Add paragraph slide
def _add_paragraph_slide(slide, paragraphs, tc):

    body_text = _rgb(tc["body_text"])
    accent = _rgb(tc["accent"])

    n = len(paragraphs)
    fsize = max(11, min(14, int(40 / max(n, 1))))

    txBox = slide.shapes.add_textbox(
        MARGIN_X,
        CONTENT_TOP,
        CONTENT_W,
        CONTENT_H
    )

    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    tf.margin_left = Inches(0.1)
    tf.margin_right = Inches(0.1)
    tf.margin_top = Inches(0.05)
    tf.margin_bottom = Inches(0.05)

    for i, para in enumerate(paragraphs):

        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

        run = p.add_run()
        run.text = para.strip()
        run.font.color.rgb = body_text
        run.font.size = Pt(fsize)
        run.font.name = "Calibri"

        p.space_before = Pt(6)
        p.space_after = Pt(6)


# Add table slide
def _add_table_slide(slide, table_data, tc):

    primary = _rgb(tc["primary"])
    prim_text = _rgb(tc["primary_text"])
    slide_bg = _rgb(tc["slide_bg"])
    body_text = _rgb(tc["body_text"])

    headers = table_data.get("headers", [])
    content = table_data.get("content", [])
    cols = max(table_data.get("columns", len(headers)), len(headers))
    data_rows = len(content)
    total_rows = data_rows + 1

    tbl_top = CONTENT_TOP
    tbl_h = min(CONTENT_H, Inches(0.45 * total_rows + 0.3))

    table_shape = slide.shapes.add_table(
        total_rows,
        cols,
        MARGIN_X,
        tbl_top,
        CONTENT_W,
        tbl_h
    )

    tbl = table_shape.table

    col_w = int(CONTENT_W / cols)

    for c in range(cols):
        tbl.columns[c].width = col_w

    for c, h in enumerate(headers[:cols]):

        cell = tbl.cell(0, c)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = primary

        p = cell.text_frame.paragraphs[0]
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = prim_text
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        cell.text_frame.margin_left = Inches(0.05)
        cell.text_frame.margin_right = Inches(0.05)
        cell.text_frame.margin_top = Inches(0.04)
        cell.text_frame.margin_bottom = Inches(0.04)

    r0, g0, b0 = tc["slide_bg"]
    alt_color = RGBColor(min(r0+20,255), min(g0+15,255), min(b0+15,255))

    for r_idx, row_data in enumerate(content):

        bg = alt_color if r_idx % 2 == 1 else None

        for c_idx in range(cols):

            cell = tbl.cell(r_idx + 1, c_idx)
            ctext = row_data[c_idx] if c_idx < len(row_data) else ""

            cell.text = str(ctext)

            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg

            p = cell.text_frame.paragraphs[0]
            p.font.size = Pt(10)
            p.font.color.rgb = body_text
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.LEFT

            cell.text_frame.margin_left = Inches(0.05)
            cell.text_frame.margin_right = Inches(0.05)
            cell.text_frame.margin_top = Inches(0.03)
            cell.text_frame.margin_bottom = Inches(0.03)


# Build PPT from refined JSON
def create_styled_table_slide(final_json, output_file_name):

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    tc = final_json.get("theme_colors", {
        "primary": [30, 80, 150],
        "primary_text": [255, 255, 255],
        "accent": [0, 180, 160],
        "slide_bg": [245, 247, 250],
        "body_text": [40, 40, 40]
    })

    slide_bg_rgb = _rgb(tc["slide_bg"])

    blank_layout = prs.slide_layouts[6]

    for slide_json in final_json["slides"]:

        slide = prs.slides.add_slide(blank_layout)

        _set_slide_bg(slide, slide_bg_rgb)

        meta = slide_json["slide_metadata"]
        layout = meta["layout_type"]

        _add_title(slide, meta["title"], tc)

        if layout == "table" and slide_json.get("table_data"):
            _add_table_slide(slide, slide_json["table_data"], tc)

        elif layout == "bullet":
            items = slide_json.get("text_content", [])
            _add_bullet_slide(slide, items, tc)

        else:
            paras = slide_json.get("text_content", [])
            _add_paragraph_slide(slide, paras, tc)

    prs.save(output_file_name)
    print(f"✅  Saved: {output_file_name}")


# Pipeline execution
def final_function(path, output_file_name, api_key):

    slides = extract_shapes(path)

    slides_json = []

    for slide_shapes in slides:

        grid = reconstruct_grid(slide_shapes)

        slide_json = build_slide_json(grid)

        slides_json.append(slide_json)

    # print("slide_json", slides_json)

    final_json = refine_json_with_gemini(slides_json, api_key)

    # print("final_json\n", final_json)

    create_styled_table_slide(final_json, output_file_name)


# Script entry
if __name__ == "__main__":

    api_key = "your_api_key"

    final_function(
        "input.pptx",
        "output.pptx",
        api_key
    )